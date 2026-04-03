"""
Use the `geo` environment
"""

#%%### Imports
import os
import sys
import site
import io
import subprocess as sp
import platform
import numpy as np
import pandas as pd
from pathlib import Path
import matplotlib as mpl
import matplotlib.pyplot as plt
import matplotlib.patheffects as pe
import geopandas as gpd
import cmocean
import folium
import mapclassify
import pyproj
import pptx
import pptx.util

### Local imports
repo_path = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
sys.path.append(repo_path)
import preprocessing

## https://github.nrel.gov/ReEDS/ReEDS-2.0
reeds_path = os.path.expanduser('~/github/ReEDS-2.0')
site.addsitedir(reeds_path)
import reeds
from reeds import plots
## https://github.nrel.gov/cobika/DLR
dlrpath = os.path.expanduser('~/github/DLR')
site.addsitedir(dlrpath)
import dlr.helpers

pyproj.network.set_network_enabled(False)
pd.options.display.max_columns = 200
plots.plotparams()

projpath = os.path.abspath(os.path.join(os.path.dirname(__file__),'..'))
maps_path = Path('~/Projects/Data/shapefiles').expanduser()

#%%### User inputs
zonepath = os.path.join(projpath, 'zones', 'z134_20030521')
# zonepath = os.path.join(projpath, 'zones', 'z132_20250313')
# zonepath = os.path.join(projpath, 'zones', '20250123_wjc_102_Zones')
# zonepath = os.path.join(projpath, 'zones', '20251028_wjc_108_Zones')
# zonepath = os.path.join(projpath, 'zones', '20250123_pb')
# zonepath = os.path.join(projpath, 'zones', '20250124_pb')
# zonepath = os.path.join(projpath, 'zones', '20250203_pb')
# zonepath = os.path.join(projpath, 'zones', '20250221_pb')
# zonepath = os.path.join(projpath, 'zones', '20250228_pb')
# zonepath = os.path.join(projpath, 'zones', '20250306_pb')
# zonepath = os.path.join(projpath, 'zones', '20250307_pb')
# zonepath = os.path.join(projpath, 'zones', '20250307_tm')
# zonepath = os.path.join(projpath, 'zones', '20251013_pb')
# zonepath = os.path.join(projpath, 'zones', '20251104_pb')
# zonepath = os.path.join(projpath, 'zones', '20251114_pb')
# zonepath = os.path.join(projpath, 'zones', '20251202_pb')
# zonepath = os.path.join(projpath, 'zones', 'zPJMcounty_20251119')
# zonepath = os.path.join(projpath, 'zones', 'zUTcounty_20251119')
# zonepath = os.path.join(projpath, 'zones', 'z48_state')
# zonepath = os.path.join(projpath, 'zones', 'z54_20220624')
# zonepath = os.path.join(projpath, 'zones', 'z69_20220624')
# zonepath = str(Path(projpath, 'zones', 'z90_20260206'))
# zonepath = str(Path(projpath, 'zones', 'z3109_20260216'))
# zonepath = str(Path(projpath, 'zones', 'z90_20260216'))
# zonepath = str(Path(projpath, 'zones', 'z153_20260223'))
# zonepath = str(Path(projpath, 'zones', 'z1259_20260223'))
# zonepath = str(Path(projpath, 'zones', 'z2975_20260223'))
# zonepath = str(Path(projpath, 'zones', 'z3109_20260223'))
# zonepath = str(Path(projpath, 'zones', 'z2972_20260303'))
## focus_states: list of 2-letter abbrevs (e.g. ['CA','NY']) or 'all'
focus_states = 'all'
interactive_layers = [
    'Independent_System_Operators',
    'Control_Areas',
    'Electric_Holding_Company_Areas',
    'Electric_Planning_Areas',
    'Electric_Retail_Service_Territories',
    'NERC_Regions',
    'NERC_Reliability_Coordinators',
    'RTO_Regions',
    'PJM_Zones_public',
    'cities',
    'urban_areas',
    'mountains',
    'transmission',
]

enforce_upper_case = False
interactive = False
# simplify = {'county':1000, 'urban':2000}
simplify = 1000
cmap = cmocean.cm.rain
focus_transmission = False
crs = 'ESRI:102008'
GSw_EFS1_AllYearLoad = 'EER_IRAlow'
loadyear = 2040
latlon_bounds = {'latmin':20, 'latmax':50, 'lonmin':-130, 'lonmax':-60}

layer_props = {
    'Independent_System_Operators': {'tooltip':'NAME'},
    'Control_Areas': {'tooltip':'NAME'},
    'Electric_Holding_Company_Areas': {'tooltip':'NAME'},
    'Electric_Planning_Areas': {'tooltip':'NAME'},
    'Electric_Retail_Service_Territories': {'tooltip':'NAME'},
    'NERC_Regions': {'tooltip':'SUBNAME'},
    'NERC_Reliability_Coordinators': {'tooltip':'NAME'},
    'RTO_Regions': {'tooltip':'Unique_ID'},
    'PJM_Zones_public': {'tooltip':'PLANNING_Z'},
    'continental_divide': {'color':'green'},
    'urban_areas': {'tooltip':'NAME20', 'color':'black'},
    'cities': {
        'tooltip': ['city','pop_2023'], 'color':'red',
        'marker_kwds': {'radius':3},
    },
    'mountains': {'tooltip':'label', 'color':'gray'},
    'transmission': {
        'layercol': 'rep_voltage',
        # 'layercolors': {230:'C8', 345:'C2', 500:'C0', 765:'C4'},
        'layercolors': {115: '0.8', 230:'0.6', 345:'0.4', 500:'0.2', 765:'0.0'},
        'tooltip': ['rep_voltage','hifld_id'],
        'weight':1.0,
    },
}


#%%### Functions
def get_counties(
    crs='ESRI:102008',
    projpath=projpath,
):
    """
    """
    ### Load and convert to CRS
    dfcounty = (
        reeds.spatial.get_map('county', source='tiger', crs=crs).reset_index()
        .astype({'INTPTLAT':float, 'INTPTLON':float})
        .rename(columns={'GEOID':'FIPS'})
        .set_index('FIPS')
    )

    ### Subset to contiguous US
    statefips = pd.read_csv(
        os.path.join(projpath,'zones','statefips.csv'),
        comment='#', dtype={'STATEFP':str},
    ).rename(columns={'STATEFP':'FIPS'}).set_index('FIPS')
    dfcounty['STATE'] = dfcounty.STATEFP.map(statefips.USPS)
    contiguous_fips = statefips.loc[~statefips.USPS.isin(['AK','HI'])]
    ## Lower 48 + Washington DC
    assert len(contiguous_fips) == 49
    dfcounty = dfcounty.loc[dfcounty.STATEFP.isin(contiguous_fips.index)].copy()

    return dfcounty


def get_urban(
    url='https://www2.census.gov/geo/tiger/TIGER2024/UAC20/tl_2024_us_uac20.zip',
    simplify=1000,
    crs='ESRI:102008',
    projpath=projpath,
):
    """
    Args:
        url (str): URL for data.
            Parent site: https://www.census.gov/cgi-bin/geo/shapefiles/index.php
        simplify (float [m]): Amount by which to simplify the geometry
    """
    ### Download and unzip the file
    local_path = os.path.join(projpath, '.cache', os.path.basename(url).replace('.zip',''))
    if not os.path.exists(local_path):
        reeds.remote.download(url, local_path)

    ### Load and convert to CRS
    urban = (
        gpd.read_file(local_path)
        .to_crs(crs)
        .astype({'INTPTLAT20':float, 'INTPTLON20':float})
    )

    ### Subset to contiguous US
    urban = urban.loc[
        (urban.INTPTLAT20 >= latlon_bounds['latmin'])
        & (urban.INTPTLAT20 <= latlon_bounds['latmax'])
        & (urban.INTPTLON20 <= latlon_bounds['lonmax'])
        & (urban.INTPTLON20 >= latlon_bounds['lonmin'])
    ].copy()

    ### Simplify the geometry to speed up mapping
    urban['geometry'] = urban.simplify(simplify, preserve_topology=True).buffer(0.)

    return urban


def get_continental_divide(crs='ESRI:102008'):
    """
    https://earthworks.stanford.edu/catalog/stanford-pw312bv3382
    """
    import requests
    # Accept user input for the GeoServer URL and the typeName of the layer
    geoserver_url = "https://geowebservices.stanford.edu/geoserver/wfs?"
    type_name = "druid:pw312bv3382"

    # Define the parameters for the WFS service
    params = {
        'service': 'WFS',
        'version': '2.0.0',
        'request': 'GetFeature',
        'typeName': type_name,
        'outputFormat': 'application/json'
    }

    # Send a GET request to the GeoServer WFS service
    response = requests.get(geoserver_url, params=params)

    # Load the features from the WFS service into a GeoDataFrame
    gdf = gpd.read_file(response.text).to_crs(crs)

    return gdf


def get_load_county(
    GSw_EFS1_AllYearLoad='EER_IRAlow',
    year=2040,
    GSw_LoadAllocationMethod='state_lpf',
):
    """
    """
    ### Get load
    assert year in [2021, 2025, 2030, 2035, 2040, 2045, 2050]
    assert GSw_LoadAllocationMethod in ['state_lpf']
    load_state = (
        reeds.io.get_load_hourly(GSw_EFS1_AllYearLoad=GSw_EFS1_AllYearLoad)
        .loc[year]
        .mean()
        .rename_axis('st')
    )

    ### Disaggregate to counties
    county_state_frac = pd.read_csv(
        Path(
            reeds_path, 'inputs', 'disaggregation', f'county_{GSw_LoadAllocationMethod}.csv',
        )
    )
    county_state_frac.FIPS = county_state_frac.FIPS.str.strip('p')
    county_state_frac = county_state_frac.set_index('FIPS').squeeze(1).rename('state_frac')

    fips2state = pd.read_csv(
        Path(reeds_path, 'inputs', 'county2zone.csv'),
        dtype=str,
    ).rename(columns={'state':'st'})

    load_mean_county = (
        load_state.rename('MW').reset_index()
        .merge(fips2state[['FIPS','st']], on='st', how='outer')
        .set_index('FIPS').MW
        * county_state_frac
        / 1e3
    ).rename('GW')

    if (1 - load_state.sum() / load_mean_county.sum()) > 0.05:
        raise ValueError('Something went wrong with load disaggregation')

    return load_mean_county


def add_to_pptx(
        title=None, file=None, left=0, top=0.62, width=13.33, height=None,
        verbose=1, slide=None,
    ):
    """
    Add current matplotlib figure (or file if specified) to new powerpoint slide
    TODO: Make this function importable from ReEDS repo
    """
    if not file:
        image = io.BytesIO()
        plt.savefig(image, format='png')
    else:
        image = file
        if not os.path.exists(image):
            raise FileNotFoundError(image)

    if slide is None:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.title.text = title
    slide.shapes.add_picture(
        image,
        left=(None if left is None else pptx.util.Inches(left)),
        top=(None if top is None else pptx.util.Inches(top)),
        width=(None if width is None else pptx.util.Inches(width)),
        height=(None if height is None else pptx.util.Inches(height)),
    )
    if verbose:
        print(title)
    return slide


#%%### Inferred inputs
savepath = os.path.join(zonepath, 'plots')
savename_html = f"maps-{len(interactive_layers)}.html"
savename_pptx = f'maps-{os.path.basename(zonepath)}.pptx'
os.makedirs(savepath, exist_ok=True)


# #%% Write for TSC
# tscpath = Path('~/github/TSC').expanduser()
# sys.path.append(str(tscpath))
# import tsc
# tsc.preprocessing.write_zones(
#     path_county2zone=Path(zonepath, 'county2zone.csv'),
#     path_hierarchy=Path(zonepath, 'hierarchy.csv'),
#     zonesname=Path(zonepath).name,
# )
# print('Zones written for tsc')


#%%### ReEDS maps (only for country and state borders)
dfmap = reeds.io.get_dfmap()
country = dfmap['country'].squeeze().geometry

#%% State FIPS codes
### https://www.census.gov/library/reference/code-lists/ansi.html#state
statefips = pd.read_csv(os.path.join(projpath,'zones','statefips.csv'), comment='#')
statefp2state = statefips.set_index('STATEFP').state

#%% Counties
_dfcounty = get_counties()
## Subset to land area
_dfcounty.geometry = _dfcounty.intersection(country)

#%% User input: county-to-zone map
county2zone = pd.read_csv(
    os.path.join(zonepath, 'county2zone.csv'),
    dtype={'FIPS':str},
    comment='#',
).dropna().rename(columns={'ba':'r'}).set_index('FIPS')
if (county2zone.r.str.upper() != county2zone.r).any() and enforce_upper_case:
    raise ValueError('Zone names should all be upper case')
numzones = len(county2zone.r.unique())
## Apply it
_dfcounty['r'] = county2zone['r']

#%% Create zone map from user input
dfzones = _dfcounty.copy()
dfzones = dfzones.dissolve('r')[['STATEFP','STATE','geometry']].copy()
dfzones['name'] = dfzones.index.copy()
## Get node location
dfzones['node'] = dfzones.geometry.map(preprocessing.spatial.get_node)

#%% Write zones for ReEDS
dfwrite = dfzones[['node','geometry']].copy()
latlon = (
    dfwrite.drop(columns='geometry').rename(columns={'node':'geometry'})
    .set_crs(dfzones.crs).to_crs('EPSG:4326')
)
latlon['latitude'] = latlon.geometry.y
latlon['longitude'] = latlon.geometry.x
dfwrite['node_latitude'] = latlon.latitude
dfwrite['node_longitude'] = latlon.longitude
## Sort if possible
try:
    dfwrite['order'] = dfwrite.index.map(lambda x: int(x.lstrip('p')))
    dfwrite = dfwrite.sort_values('order').drop(columns='order')
except ValueError:
    dfwrite = dfwrite.sort_index()
dfwrite[['node_latitude','node_longitude','geometry']].to_file(
    os.path.join(zonepath, 'modelzones.gpkg')
)


#%% Simplify counties from now on to speed up plots
dfcounty = _dfcounty.copy()
dfcounty['geometry'] = dfcounty.simplify(simplify, preserve_topology=True).buffer(0.)
# dfcounty['geometry'] = dfcounty.simplify_coverage(simplify).buffer(0.)


#%% Get load for counties and zones
load_mean_county = get_load_county(
    GSw_EFS1_AllYearLoad=GSw_EFS1_AllYearLoad, year=loadyear)

## Zonal mean load
load_mean_zone = county2zone.copy()
load_mean_zone['GW'] = load_mean_county
load_mean_zone = load_mean_zone.groupby('r').GW.sum()


#%% Get current generation capacity
dfeia = pd.read_csv(
    os.path.join(
        reeds_path, 'inputs', 'capacity_exogenous',
        'ReEDS_generator_database_final_EIA-NEMS.csv'),
)
dfeia['FIPS'] = dfeia['FIPS'].str.strip('p')
dfeia['zone'] = dfeia['FIPS'].map(county2zone.r)
dfeia = dfeia.loc[
    (dfeia.RetireYear >= 2026)
    & (dfeia.StartYear <= 2025)
].copy()
cap_county = dfeia.groupby('FIPS').summer_power_capacity_MW.sum() / 1e3

cap_zone = county2zone.copy()
cap_zone['GW'] = cap_county
cap_zone = cap_zone.groupby('r').GW.sum()


#%% Contingency stuff
largest_unit = (
    dfeia.sort_values('summer_power_capacity_MW', ascending=False)
    .groupby('zone').head(1)
    .set_index('zone')
    .sort_index()
)
contingency_pct = largest_unit.summer_power_capacity_MW / (load_mean_zone * 1e3) * 100
cap_meanload_pct = dfeia.groupby('zone').summer_power_capacity_MW.sum() / (load_mean_zone * 1e3) * 100
capminuslargest_meanload_pct = (
    (dfeia.groupby('zone').summer_power_capacity_MW.sum() - largest_unit.summer_power_capacity_MW)
    / (load_mean_zone * 1e3) * 100
)


#%%### Get optional maps
map_layers = {}

#%% Transmission: HIFLD
if 'transmission' in interactive_layers:
    try:
        map_layers['transmission'] = (
            dlr.helpers.get_lines(
                min_kv=115, max_miles=np.inf,
                # remove_dc=False, remove_underground=False,
            )
            .astype({'rep_voltage':int})
        )
        map_layers['transmission']['hifld_id'] = map_layers['transmission'].index
        map_layers['transmission'].geometry = map_layers['transmission'].simplify(simplify/2)
    except FileNotFoundError as err:
        print(err)

#%% Cities by population
if 'cities' in interactive_layers:
    try:
        cities = pd.read_csv(os.path.join(projpath,'zones','population_city.csv'), comment='#')
        cities = cities.loc[~cities.state.isin(['AK','HI'])].copy()
        cities = plots.df2gdf(cities)
        map_layers['cities'] = cities
    except FileNotFoundError as err:
        print(err)

#%% Urban areas
if 'urban_areas' in interactive_layers:
    try:
        map_layers['urban_areas'] = get_urban()
    except FileNotFoundError as err:
        print(err)

#%% ISOs / RTOs
grid_areas = [
    'Control_Areas',
    'Electric_Holding_Company_Areas',
    'Electric_Planning_Areas',
    'Electric_Retail_Service_Territories',
    'Independent_System_Operators',
    'NERC_Regions',
    'NERC_Reliability_Coordinators',
    'RTO_Regions',
    'PJM_Zones_public',
]
for grid_area in interactive_layers:
    if grid_area in grid_areas:
        try:
            for ftype in ['.gpkg', '.geojson', '']:
                fpath = os.path.join(maps_path, f'{grid_area}{ftype}')
                if os.path.exists(fpath):
                    print(os.path.basename(fpath))
                    map_layers[grid_area] = gpd.read_file(fpath).to_crs(crs)
                    continue
            ## Simplify and subset to contiguous US
            if grid_area not in map_layers:
                print(f'{fpath} not found')
                continue
            map_layers[grid_area].geometry = (
                map_layers[grid_area]
                .simplify(simplify).buffer(0.)
                .intersection(country)
            )
        except FileNotFoundError as err:
            print(err)

#%% Other features
if 'mountains' in interactive_layers:
    fpath = os.path.join(maps_path, 'mountains.gpkg')
    if not os.path.exists(fpath):
        print(f'{fpath} not found')
    else:
        map_layers['mountains'] = gpd.read_file(fpath)
        map_layers['mountains'] = (
            map_layers['mountains']
            .loc[map_layers['mountains']['class'].isin(['high','high_scattered'])]
            .dissolve()
            .assign(label='mountains')
            .to_crs(crs)
        )
        # map_layers['mountains'].geometry = map_layers['mountains'].simplify(simplify)

if 'continental_divide' in interactive_layers:
    try:
        map_layers['continental_divide'] = get_continental_divide()
        map_layers['continental_divide'].geometry = (
            map_layers['continental_divide'].simplify(simplify))
    except FileNotFoundError as err:
        print(err)


#%%### Static maps
### Set up powerpoint file
prs = pptx.Presentation(os.path.join(reeds_path,'postprocessing','template.pptx'))
blank_slide_layout = prs.slide_layouts[3]


#%% Differentiated zone map
## Get colors
try:
    colors = 'C' + mapclassify.greedy(dfzones, strategy='smallest_last').astype(str)
except Exception as err:
    print(err)
    colors = {}
    for state in dfzones.STATE.unique():
        zones = dfzones.loc[dfzones.STATE==state, 'name']
        for i, zone in enumerate(zones):
            colors[zone] = f"C{i%10}"

### Plot it
for label_zone in [True, False]:
    plt.close()
    f,ax = plt.subplots(figsize=(10,6))
    # dfcounty.plot(ax=ax, facecolor='none', edgecolor='C7', lw=0.1, alpha=0.4, zorder=1e5)
    dfzones.plot(ax=ax, facecolor='none', edgecolor='k', lw=0.3, zorder=1e6)
    dfmap['st'].plot(ax=ax, facecolor='none', edgecolor='k', lw=0.6, zorder=1e7)
    for r, row in dfzones.iterrows():
        dfzones.loc[[r]].plot(ax=ax, color=colors[r], alpha=0.4, lw=0, zorder=1e4)
        if label_zone:
            ax.annotate(
                r, (row.geometry.centroid.x, row.geometry.centroid.y),
                ha='center', va='center', size=7, color='k', zorder=1e8,
                path_effects=[pe.withStroke(linewidth=1.5, foreground='w', alpha=0.9)]
            )
    ax.axis('off')

    title = f'{numzones} zones ({os.path.basename(zonepath)})'
    slide = add_to_pptx(title, width=None, height=6.88)
    if interactive:
        plt.show()


#%% County load
dfplot = dfcounty.copy()
dfplot['GW'] = load_mean_county
vmax = dfplot.GW.max()

plt.close()
f,ax = plt.subplots(figsize=(12,9))
dfmap['st'].plot(ax=ax, facecolor='none', edgecolor='k', lw=0.9, zorder=3e7)
dfplot.plot(ax=ax, facecolor='none', edgecolor='C7', lw=0.2, zorder=1e7)
dfplot.plot(ax=ax, column='GW', cmap=plt.cm.gist_heat_r, vmin=0, vmax=vmax)
dfzones.plot(ax=ax, facecolor='none', edgecolor='k', lw=0.6, zorder=2e7)
ax.axis('off')
for r, row in dfplot.loc[dfplot.GW >= 1].iterrows():
    ax.annotate(
        f'{row.GW:.0f}', (row.geometry.centroid.x, row.geometry.centroid.y),
        ha='center', va='center', color=('k' if row.GW <= vmax*0.5 else 'w'), fontsize=6)
ax.set_title(
    f'{loadyear} average demand, {GSw_EFS1_AllYearLoad} [GW]',
    va='bottom', y=0.95, fontsize='x-large')

title = 'Mean demand by county'
slide = add_to_pptx(title, width=None, height=6.88)
if interactive:
    plt.show()


#%% Zonal load
plotsettings = [
    {
        'ylabel':'Mean load [GW]',
        'slidetitle':'Mean demand by zone',
        'df':load_mean_zone,
        'title':f'{loadyear} average demand, {GSw_EFS1_AllYearLoad} [GW]',
        'yscale':'log',
        'cmap':cmocean.cm.rain,
    },
    {
        'ylabel':'Gen cap [GW]',
        'slidetitle':'Generation capacity by zone',
        'df':cap_zone,
        'title':'2025 generation capacity [GW]',
        'yscale':'log',
        'cmap':cmocean.cm.rain,
    },
    # {
    #     'ylabel':'Contingency [%]',
    #     'slidetitle':'Largest contingency by zone',
    #     'df':contingency_pct,
    #     'title':f'Largest unit capacity / {loadyear} mean demand ({GSw_EFS1_AllYearLoad})',
    #     'yscale':'linear',
    #     'cmap':cmocean.cm.rain,
    # },
    {
        'ylabel':'(g-0) cap / load [%]',
        'slidetitle':'Capacity (g-0) / mean load',
        'df':cap_meanload_pct,
        'title':f'Generation capacity (g-0) / {loadyear} mean demand ({GSw_EFS1_AllYearLoad})',
        'yscale':'linear',
        'cmap':plt.cm.coolwarm_r,
        'vmax': 200,
    },
    {
        'ylabel':'(g-1) cap / load [%]',
        'slidetitle':'Capacity (g-1) / mean load',
        'df':capminuslargest_meanload_pct,
        'title': (
            f"Generation capacity – largest contingency (g-1) / "
            f"{loadyear} mean demand ({GSw_EFS1_AllYearLoad})"),
        'yscale':'linear',
        'cmap':plt.cm.coolwarm_r,
        'vmax': 200,
    },
]

for plotdict in plotsettings:
    dfplot = dfzones.copy()
    dfplot['val'] = plotdict['df']

    vmax = plotdict.get('vmax', dfplot.val.max())
    plt.close()
    f,ax = plt.subplots(figsize=(12,9))
    dfplot.plot(ax=ax, facecolor='none', edgecolor='k', lw=0.4, zorder=1e7)
    dfplot.plot(ax=ax, column='val', cmap=plotdict['cmap'], vmin=0, vmax=vmax)
    dfcounty.plot(ax=ax, facecolor='none', edgecolor='C7', lw=0.1, alpha=0.4, zorder=1e6)
    for r, row in dfplot.iterrows():
        if plotdict['cmap'] in [plt.cm.turbo, plt.cm.turbo_r, plt.cm.coolwarm, plt.cm.coolwarm_r]:
            textcolor = 'k' if vmax * 0.2 <= row.val <= vmax * 0.8 else '0.9'
        else:
            textcolor = 'k' if row.val <= vmax * 0.5 else '0.9'
        ax.annotate(
            f'{row.val:.1f}', (row.geometry.centroid.x, row.geometry.centroid.y),
            ha='center', va='center', color=textcolor,
            zorder=1e8,
        )
    ax.axis('off')
    ax.set_title(
        plotdict['title'],
        va='bottom', y=0.95, fontsize='x-large')

    ## Subplot: sorted load by zone
    sax = f.add_axes([0.21, 0.21, 0.2, 0.15])
    sax.plot(
        range(len(plotdict['df'])),
        plotdict['df'].sort_values(ascending=False).values,
        marker='o', markeredgewidth=0, ms=3, lw=1,
    )
    sax.set_facecolor('none')
    sax.set_ylabel(plotdict['ylabel'])
    sax.set_xlabel('Number of zones')
    sax.set_xlim(-1)
    if plotdict.get('yscale') == 'log':
        sax.set_yscale('log')
    else:
        sax.set_ylim(0)
        sax.yaxis.set_minor_locator(mpl.ticker.AutoMinorLocator(5))
    sax.xaxis.set_major_locator(mpl.ticker.MultipleLocator(20))
    sax.xaxis.set_minor_locator(mpl.ticker.MultipleLocator(5))
    plots.despine(sax)

    slide = add_to_pptx(plotdict['slidetitle'], width=None, height=6.88)
    if interactive:
        plt.show()


#%% Other map layers
for key, dflayer in map_layers.items():
    if key in ['transmission', 'continental_divide']:
        continue
    ## Get colors
    plt.close()
    f,ax = plt.subplots(figsize=(12,9))
    ## Layer
    if len(dflayer) <= 200:
        if len(dflayer) <= 21:
            layercolors = plots.rainbowmapper(dflayer.index, categorical=True)
        else:
            layercolors = (
                'C' + mapclassify.greedy(dflayer, strategy='smallest_last').astype(str)
            )
        for i, row in dflayer.iterrows():
            dflayer.loc[[i]].plot(
                ax=ax, facecolor=layercolors[i], edgecolor='none', alpha=0.7,
            )
    else:
        dflayer.plot(ax=ax, facecolor='C3', edgecolor='none', alpha=0.7)
    ## Rest of plot
    ax.set_xlim(*ax.get_xlim())
    ax.set_ylim(*ax.get_ylim())
    dfcounty.plot(ax=ax, facecolor='none', edgecolor='0.5', lw=0.1, zorder=1e7)
    dfzones.plot(ax=ax, facecolor='none', edgecolor='k', lw=0.8, zorder=2e7)
    ax.axis('off')

    title = key
    slide = add_to_pptx(title, width=None, height=6.88)
    if interactive:
        plt.show()


#%%### Hierarchy levels
colors_hierarchy = {}
colors_hierarchy['transgrp'] = {
    'CAISO':plt.cm.tab20c(6),
    'NorthernGrid_West':plt.cm.tab20c(1),
    'NorthernGrid_East':plt.cm.tab20c(2),
    'NorthernGrid_South':plt.cm.tab20c(3),
    'WestConnect_North':plt.cm.tab20c(9),
    'WestConnect_South':plt.cm.tab20c(10),
    'SPP_North':plt.cm.tab20c(5),
    'SPP_South':plt.cm.tab20c(6),
    'MISO_North':plt.cm.tab20c(9),
    'MISO_Central':plt.cm.tab20c(10),
    'MISO_South':plt.cm.tab20c(11),
    'ERCOT':plt.cm.tab20c(2),
    'PJM_West':plt.cm.tab20c(5),
    'PJM_East':plt.cm.tab20c(6),
    'SERTP':plt.cm.tab20c(2),
    'FRCC':plt.cm.tab20c(10),
    'NYISO':plt.cm.tab20c(2),
    'ISONE':plt.cm.tab20c(10),
}
offset = {}
offset['transgrp'] = {
    'MISO_Central': (-1e5,-1e5),
    'CAISO': (-0.2e5,0),
}

if os.path.exists(os.path.join(zonepath, 'hierarchy.csv')):
    hierarchy = pd.read_csv(
        os.path.join(zonepath, 'hierarchy.csv')
    ).rename(columns={'ba':'r'}).set_index('r')

    for level in hierarchy:
        dfzones[level] = hierarchy[level]
        dfmap[level] = gpd.GeoDataFrame(
            dfzones.dissolve(level).buffer(0.).rename('geometry')
        )
        dfmap[level]['centroid_x'] = dfmap[level].centroid.x
        dfmap[level]['centroid_y'] = dfmap[level].centroid.y

        if level not in colors_hierarchy:
            try:
                colors_hierarchy[level] = (
                'C'
                + mapclassify.greedy(dfmap[level], strategy='smallest_last').astype(str)
            )
            except Exception as err:
                print(err)
                _colors = reeds.plots.rainbowmapper(dfmap[level].index, categorical=True)
                colors_hierarchy[level] = dict(zip(
                    dfmap[level].index,
                    dfmap[level].index.map(_colors)
                ))

        plt.close()
        f,ax = plt.subplots(figsize=(10,6))
        # dfcounty.plot(ax=ax, facecolor='none', edgecolor='C7', lw=0.1, alpha=0.4, zorder=1e5)
        dfzones.plot(ax=ax, facecolor='none', edgecolor='C7', lw=0.3, zorder=1e6)
        dfmap['st'].plot(ax=ax, facecolor='none', edgecolor='C7', lw=0.6, zorder=1e7)
        dfmap[level].plot(ax=ax, facecolor='none', edgecolor='k', lw=1.0, zorder=1e7)
        for r, row in dfmap[level].iterrows():
            dfmap[level].loc[[r]].plot(
                ax=ax, color=colors_hierarchy[level][r], alpha=0.8, lw=0, zorder=1e4)
            x = row.centroid_x + offset.get(level,{}).get(r, (0,0))[0]
            y = row.centroid_y + offset.get(level,{}).get(r, (0,0))[1]
            ax.annotate(
                r.replace('_','\n'), (x, y),
                ha='center', va='center', color='k', zorder=1e8,
                path_effects=[pe.withStroke(linewidth=1.5, foreground='w', alpha=0.9)],
                size=10, weight='bold',
            )
        ax.axis('off')

        title = f"{level} ({len(dfmap[level])})"
        slide = add_to_pptx(title, width=None, height=6.88)
        if interactive:
            plt.show()


#%%### Focus states
if focus_states == 'all':
    state_count = dfzones.groupby('STATE').name.count()
    _focus_states = state_count.loc[state_count > 1].index
elif isinstance(focus_states, str):
    _focus_states = [focus_states]
else:
    _focus_states = focus_states

for state in _focus_states:
    df = dfcounty.loc[dfcounty.STATE==state].copy()
    df['GW'] = load_mean_county
    # vmax = df.GW.max()
    vmax = 10
    if not len(df):
        print(f'No counties found for state={state}. Use two-letter abbreviation.')

    plt.close()
    f,ax = plt.subplots(figsize=(13.33,6.88))
    df.plot(ax=ax, facecolor='none', edgecolor='C7', lw=0.2, zorder=1e7)
    df.plot(ax=ax, column='GW', cmap=plt.cm.binary, vmin=0, vmax=vmax, zorder=1e6)
    for r, row in dfzones.loc[dfzones.STATE==state].iterrows():
        dfzones.loc[[r]].plot(
            ax=ax, facecolor='none', edgecolor=colors[r],
            lw=0.5, zorder=1e8,
        )
        dfzones.loc[[r]].plot(
            ax=ax, facecolor=colors[r], edgecolor='none',
            lw=0, zorder=1e8, alpha=0.05,
        )
    for i, row in df.iterrows():
        ax.annotate(
            row.NAME, (row.geometry.centroid.x, row.geometry.centroid.y),
            ha='center', va='center', fontsize=7, zorder=1e10,
            color=colors[county2zone.loc[i,'r']],
            path_effects=[pe.withStroke(linewidth=1.5, foreground='w', alpha=0.9)],
        )
    if focus_transmission:
        ax.set_xlim(*ax.get_xlim())
        ax.set_ylim(*ax.get_ylim())
        map_layers['transmission'].plot(ax=ax, color='m', lw=0.25, alpha=0.5, zorder=1e9)
    plots.addcolorbarhist(
        f=f, ax0=ax, data=df.GW.values, nbins=101,
        title=f'Mean\ncounty\ndemand\n{loadyear} [GW]',
        vmin=0, vmax=vmax, cmap=plt.cm.binary,
        cbarleft=1.1,
    )
    ax.axis('off')

    title = f'{state}: {len(df.r.unique())}'
    slide = add_to_pptx(title, width=None, height=6.88)
    if interactive:
        plt.show()


#%%### Full interactive map
alpha = 0.2
### Background
m = dfmap['st'].reset_index().explore(
    color='none',
    tooltip='st',
    style_kwds={'color':'white', 'weight':1.5},
    name='states')
### Tiles
tileses = ['CartoDB positron','CartoDB dark_matter','OpenStreetMap']
for tiles in tileses:
    folium.TileLayer(tiles).add_to(m)
### Zones
for r, row in dfzones.iterrows():
    dfzones.loc[[r]].explore(
        m=m,
        name='zones',
        tooltip=['name'],
        color=mpl.colors.to_hex(colors[r]),
        style_kwds={'weight':0.5, 'fillOpacity':alpha},
    )
### Counties
dfcounty.reset_index().explore(
    m=m,
    name='counties',
    tooltip=['r','FIPS','NAME','STATE'],
    color='white',
    style_kwds={'weight':0.25, 'fillOpacity':alpha/2},
)
### Other layers
for label, layer in map_layers.items():
    this_layer_props = layer_props.get(label,{})
    if 'layercol' in this_layer_props:
        for level, color in this_layer_props.get('layercolors',{}).items():
            df = layer.loc[layer[this_layer_props['layercol']] == level]
            df.explore(
                m=m,
                name=f"{label} {level}",
                tooltip=this_layer_props.get('tooltip',None),
                color=mpl.colors.to_hex(color),
                weight=this_layer_props.get('weight',None),
            )
    else:
        layer.explore(
            m=m,
            name=label,
            tooltip=this_layer_props.get('tooltip',None),
            color=mpl.colors.to_hex(this_layer_props.get('color', 'C7')),
            style_kwds={'fillOpacity':alpha},
            marker_kwds=this_layer_props.get('marker_kwds',{}),
        )

folium.LayerControl().add_to(m)
m

#%% Save it
m.save(os.path.join(savepath,savename_html))
fpath_ppt = os.path.join(savepath,savename_pptx)
prs.save(fpath_ppt)
print(fpath_ppt)

### Open it
if sys.platform == 'darwin':
    sp.run(f"open '{os.path.join(savepath,savename_html)}'", shell=True)
    sp.run(f"open '{fpath_ppt}'", shell=True)
elif platform.system() == 'Windows':
    sp.run(f'"{savename_html}"', shell=True)
    sp.run(f'"{fpath_ppt}"', shell=True)
